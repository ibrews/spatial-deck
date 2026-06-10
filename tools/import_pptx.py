"""Import a .pptx deck into Spatial Deck's SECTIONS array.

Pipeline:
  1. python-pptx extracts slide text, notes, and embedded images (deterministic).
  2. llama3.1:8b on Sam rewrites each slide's text into the
     {title, subtitle, bullets[]} shape Spatial Deck expects (single-pass).
  3. Images are copied to media/import-<hash>/ with deterministic names.
  4. Output JSON is a SECTIONS-compatible chapter with one `cases` entry per
     slide, ready to be spliced by tools/merge_sections.py.

Usage:
    python3 tools/import_pptx.py deck.pptx
    python3 tools/import_pptx.py deck.pptx --out imported.json --chapter "DECK"
    python3 tools/import_pptx.py deck.pptx --limit 5   # quick test

No LLM output is trusted verbatim — we validate types and fall back to the
raw extracted text if a slide fails to normalize.
"""
from __future__ import annotations
import argparse
import hashlib
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, str(Path(__file__).parent))
from fleet_client import call_json, ENDPOINTS  # noqa: E402

REPO_ROOT = Path(__file__).resolve().parent.parent
MEDIA_ROOT = REPO_ROOT / "media"

PROMPT = """You are normalizing one slide from a PowerPoint deck into a fixed JSON schema for a presentation framework.

Slide title (raw): {title!r}
Slide body text (raw, newline-separated):
---
{body}
---
Speaker notes (may be empty):
---
{notes}
---

Output a JSON object with EXACTLY these keys:
  title    — short punchy headline, <= 60 chars. May contain one "\\n" for a line break. Keep it close to the original title if one was given.
  subtitle — optional one-line context, <= 80 chars. Empty string if nothing fits.
  bullets  — array of 2–4 short punchy bullet strings, each <= 140 chars. Rewrite bloated sentences into tight declarative bullets. Drop filler. Preserve numbers and proper nouns exactly.

Rules:
- No markdown, no emoji unless present in original.
- Do not invent facts — if body is empty, return bullets:[] and use the title alone.
- If the original is already tight, keep it as-is.
"""


def extract_text(shape) -> list[str]:
    """Return non-empty paragraph strings from a shape, or []."""
    if not shape.has_text_frame:
        return []
    out = []
    for p in shape.text_frame.paragraphs:
        t = "".join(r.text for r in p.runs).strip()
        if t:
            out.append(t)
    return out


def slide_payload(slide) -> tuple[str, list[str], str]:
    """(title, body_lines, notes) — deterministic extraction, no LLM."""
    title = ""
    body_lines: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape == slide.shapes.title:
            title = shape.text_frame.text.strip()
        elif shape.has_text_frame:
            body_lines.extend(extract_text(shape))
    notes = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    return title, body_lines, notes


IMG_EXT = {"image/png": ".png", "image/jpeg": ".jpg", "image/jpg": ".jpg",
           "image/gif": ".gif", "image/svg+xml": ".svg", "image/webp": ".webp"}


def save_images(slide, idx: int, out_dir: Path) -> list[Path]:
    """Extract all pictures from a slide to out_dir, return saved relative paths."""
    saved = []
    pic_n = 0
    for shape in slide.shapes:
        if shape.shape_type != 13:  # MSO_SHAPE_TYPE.PICTURE
            continue
        img = shape.image
        ext = IMG_EXT.get(img.content_type, ".bin")
        digest = hashlib.sha1(img.blob).hexdigest()[:8]
        name = f"slide{idx:02d}-{pic_n}-{digest}{ext}"
        out = out_dir / name
        if not out.exists():
            out.write_bytes(img.blob)
        saved.append(out.relative_to(REPO_ROOT))
        pic_n += 1
    return saved


def normalize_with_llm(title: str, body: list[str], notes: str) -> dict:
    body_str = "\n".join(body) or "(no body text)"
    notes_str = notes or "(none)"
    data = call_json(
        "llama3.1:8b",
        PROMPT.format(title=title, body=body_str, notes=notes_str),
        endpoint=ENDPOINTS["sam"],
        required_keys=["title", "subtitle", "bullets"],
    )
    # Type validation — fall back on raw on any structural issue.
    t = data.get("title")
    if not isinstance(t, str) or not t.strip():
        raise ValueError(f"bad title: {t!r}")
    s = data.get("subtitle", "")
    if not isinstance(s, str):
        s = ""
    b = data.get("bullets", [])
    if not isinstance(b, list) or not all(isinstance(x, str) for x in b):
        raise ValueError(f"bad bullets: {b!r}")
    return {"title": t.strip(), "subtitle": s.strip(), "bullets": [x.strip() for x in b if x.strip()]}


def fallback_entry(title: str, body: list[str]) -> dict:
    return {
        "title": title or "Untitled",
        "subtitle": "",
        "bullets": body[:4] if body else [],
    }


SAFE = re.compile(r"[^a-zA-Z0-9]+")


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--out", type=Path, default=None)
    ap.add_argument("--chapter", default="IMPORTED", help="year/chapter label")
    ap.add_argument("--accent", default="teal", choices=["teal", "purple", "amber", "rose"])
    ap.add_argument("--limit", type=int, default=0, help="max slides to process (0=all)")
    ap.add_argument("--no-llm", action="store_true",
                    help="Skip LLM normalization; use raw extracted text verbatim.")
    args = ap.parse_args()

    if not args.pptx.exists():
        print(f"ERROR: {args.pptx} not found", file=sys.stderr)
        return 2

    stem = SAFE.sub("-", args.pptx.stem).strip("-").lower() or "deck"
    digest = hashlib.sha1(args.pptx.read_bytes()).hexdigest()[:8]
    media_dir = MEDIA_ROOT / f"import-{stem}-{digest}"
    media_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(args.pptx))
    slides = list(prs.slides)
    if args.limit:
        slides = slides[:args.limit]
    print(f"[pptx] {len(slides)} slides. Media → {media_dir.relative_to(REPO_ROOT)}", file=sys.stderr)

    cases = []
    for i, slide in enumerate(slides, 1):
        title, body, notes = slide_payload(slide)
        imgs = save_images(slide, i, media_dir)
        img_path = str(imgs[0]) if imgs else ""

        if args.no_llm:
            entry = fallback_entry(title, body)
        else:
            try:
                entry = normalize_with_llm(title, body, notes)
                tag = "llm"
            except Exception as e:
                print(f"  [slide {i}] LLM failed ({e}); using raw fallback", file=sys.stderr)
                entry = fallback_entry(title, body)
                tag = "raw"
            print(f"  [slide {i}/{len(slides)}] {tag}: {entry['title'][:50]}", file=sys.stderr)

        entry["img"] = img_path
        if notes:
            entry["notes"] = notes
        cases.append(entry)

    chapter = {
        "year": args.chapter,
        "accent": args.accent,
        "lesson": {
            "title": f"Imported: {args.pptx.stem}",
            "short": args.chapter,
            "tagline": f"Imported from {args.pptx.name}. Edit this tagline and the cases below to match your narrative.",
            "tags": "",
        },
        "cases": cases,
    }

    out_path = args.out or (REPO_ROOT / "tools" / f"imported-{stem}-{digest}.json")
    out_path.write_text(json.dumps(chapter, indent=2, ensure_ascii=False))
    try:
        rel = str(out_path.relative_to(REPO_ROOT))
    except ValueError:
        rel = str(out_path)
    print(f"[done] Wrote {rel} ({len(cases)} cases)", file=sys.stderr)
    print(f"[next] python3 tools/merge_sections.py {rel}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
