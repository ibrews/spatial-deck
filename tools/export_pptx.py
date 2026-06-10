"""Export the deck as a .pptx — companion to tools/import_pptx.py.

Two modes, two different honesty trade-offs:

DEFAULT (structural — editable approximation): each chapter gets a dark
"lesson" cover slide; each case gets a two-column slide (text left, image
right if present). Speaker notes land in the notes placeholder so the deck
survives re-import cleanly. This is NOT what the live deck looks like:
advanced layouts (big/placed/asset-belt/…) collapse to two columns, media
cyclers/iframes/videos are skipped, SVG images can't be embedded, and the
theme is reduced to four accent RGBs on a dark background. A fidelity
report after export lists exactly what degraded — loss is informed, never
silent.

--visual (pixel-faithful — not editable): headless Chrome captures every
visible slide via index.html?shot=N (tools/capture_slides.py — the deck
itself is the renderer, so layouts/themes/animations are exact), and each
PPTX slide is the full-bleed 16:9 screenshot with the SECTIONS speaker
notes attached to the notes placeholder. This is how Gamma/Pitch do "give
me the PowerPoint". Requires Google Chrome (or CHROME_BIN); ~1-2 min for a
full deck. The fidelity report in this mode comes from the deck's own
export manifest (cycler frozen to first item, iframe -> placeholder,
video -> frozen frame) since visual capture preserves layouts.

Usage:
    python3 tools/export_pptx.py --out /tmp/deck.pptx
    python3 tools/export_pptx.py --chapter 0 --out /tmp/ch0.pptx
    python3 tools/export_pptx.py --visual --out /tmp/deck-visual.pptx
    python3 tools/export_pptx.py --visual --slides 1-12 --out /tmp/part1.pptx
    python3 tools/export_pptx.py --out deck.pptx --json   # report as JSON
"""
from __future__ import annotations
import argparse
import json
import re
import shutil
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
from lint_deck import extract_sections  # noqa: E402

REPO_ROOT = Path(__file__).resolve().parent.parent
DEFAULT_HTML = REPO_ROOT / "index.html"

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dgm.color import RGBColor  # type: ignore  # noqa
except Exception:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
    except ImportError:
        print("ERROR: python-pptx not installed. pip3 install --user --break-system-packages python-pptx", file=sys.stderr)
        sys.exit(2)

# 16:9 default canvas is 13.33 × 7.5 inches.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ACCENT_RGB = {
    "teal":   RGBColor(0x0F, 0xB3, 0xA1),
    "purple": RGBColor(0x7C, 0x5C, 0xFF),
    "amber":  RGBColor(0xF5, 0xA5, 0x24),
    "rose":   RGBColor(0xF4, 0x3F, 0x7D),
}
BG = RGBColor(0x0B, 0x0B, 0x10)
FG = RGBColor(0xE8, 0xE8, 0xF0)
DIM = RGBColor(0xB8, 0xB8, 0xC8)

VIDEO_EXTS = (".mp4", ".webm", ".mov", ".m4v")


def _fill_bg(slide, rgb=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_textbox(slide, left, top, width, height, text, *,
                 size=18, bold=False, color=FG, line_height=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = (text or "").split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_height
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def _is_local_image(img: str) -> Path | None:
    if not img or img.startswith(("IFRAME:", "MEDIA_CYCLER", "data:")):
        return None
    p = (REPO_ROOT / img) if not img.startswith("/") else Path(img)
    try:
        if p.exists() and p.suffix.lower() in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
            return p
    except OSError:
        return None
    return None


def add_chapter_cover(prs: "Presentation", ch: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill_bg(slide)
    accent = ACCENT_RGB.get((ch.get("accent") or "").strip(), ACCENT_RGB["teal"])
    lesson = ch.get("lesson") or {}
    year = (ch.get("year") or "").strip()
    short = (lesson.get("short") or "").strip()
    header = year + ((" · " + short) if short and short != year else "")
    _add_textbox(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(0.5),
                 header, size=14, bold=True, color=accent)
    title = (lesson.get("title") or "").strip()
    _add_textbox(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(3),
                 title, size=44, bold=True, color=FG)
    tagline = (lesson.get("tagline") or "").strip()
    if tagline:
        _add_textbox(slide, Inches(0.7), Inches(5.2), Inches(12), Inches(1.6),
                     tagline, size=18, color=DIM)
    tags = (lesson.get("tags") or "").strip()
    if tags:
        _add_textbox(slide, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
                     tags, size=12, color=accent)


def add_case_slide(prs: "Presentation", c: dict, accent: "RGBColor"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    img_path = _is_local_image(c.get("img") or "")
    text_w = Inches(6.3) if img_path else Inches(12)

    _add_textbox(slide, Inches(0.7), Inches(0.7), text_w, Inches(1.2),
                 (c.get("title") or "").strip() or "(untitled)",
                 size=30, bold=True, color=FG)
    subtitle = (c.get("subtitle") or "").strip()
    if subtitle:
        _add_textbox(slide, Inches(0.7), Inches(1.9), text_w, Inches(0.9),
                     subtitle, size=16, color=DIM)

    bullets = [b for b in (c.get("bullets") or []) if isinstance(b, str) and b.strip()]
    if bullets:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), text_w, Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.3
            run = p.add_run()
            run.text = "• " + b.strip()
            run.font.size = Pt(16)
            run.font.color.rgb = FG

    if img_path:
        try:
            slide.shapes.add_picture(str(img_path), Inches(7.3), Inches(0.7),
                                     width=Inches(5.3), height=Inches(6.1))
        except Exception as e:
            print(f"    [pptx] image skipped ({img_path.name}): {e}", file=sys.stderr)

    notes = (c.get("notes") or "").strip()
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    # Accent bar on the left edge.
    from pptx.shapes.autoshape import Shape  # noqa: F401
    bar = slide.shapes.add_shape(1, Inches(0.3), Inches(0.7), Inches(0.08), Inches(6.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def build(sections: list[dict], out_path: Path) -> tuple[int, int]:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    n_cases = 0
    for ch in sections:
        add_chapter_cover(prs, ch)
        accent = ACCENT_RGB.get((ch.get("accent") or "").strip(), ACCENT_RGB["teal"])
        for c in ch.get("cases") or []:
            add_case_slide(prs, c, accent)
            n_cases += 1
    prs.save(str(out_path))
    return len(sections), n_cases


# ── Fidelity report (structural mode — derived from SECTIONS, no Chrome) ────

def _extra_media(c: dict) -> tuple[int, int]:
    """(item_count, video_count) across placedImages/placedItems/beltItems."""
    srcs: list[str] = []
    for it in c.get("placedImages") or []:
        if isinstance(it, (list, tuple)) and it:
            srcs.append(str(it[0]))
    for it in (c.get("placedItems") or []) + (c.get("beltItems") or []):
        if isinstance(it, dict):
            srcs.append(str(it.get("src") or it.get("type") or ""))
    vids = sum(1 for s in srcs if s.lower().endswith(VIDEO_EXTS) or s == "mp4")
    return len(srcs), vids


def structural_fidelity(sections: list[dict]) -> list[dict]:
    """Per-case degradations the two-column approximation introduces."""
    drops = []
    for si, ch in enumerate(sections):
        for ci, c in enumerate(ch.get("cases") or []):
            notes = []
            layout = (c.get("layout") or "").strip()
            if layout:
                notes.append(f"layout '{layout}' collapsed to two-column text+image")
            big_fields = [k for k in ("bigText", "bigCaption") if (c.get(k) or "").strip()]
            if big_fields:
                notes.append(f"fields dropped: {', '.join(big_fields)}")
            n_items, n_vids = _extra_media(c)
            if n_items:
                vid_note = f", {n_vids} of them video(s)" if n_vids else ""
                notes.append(f"{n_items} placed/belt item(s) dropped{vid_note}")
            img = (c.get("img") or "").strip()
            if img.startswith("MEDIA_CYCLER"):
                notes.append("media cycler -> skipped (no media embedded)")
            elif img.startswith("IFRAME:"):
                notes.append(f"iframe -> skipped ({img[7:]})")
            elif img.lower().endswith(VIDEO_EXTS):
                notes.append("video -> dropped (PPTX structural export embeds images only)")
            elif img and not _is_local_image(img):
                notes.append(f"image not embeddable -> skipped ({img})")
            if notes:
                title = (c.get("title") or "(untitled)").replace("\n", " ").strip()
                drops.append({"slide": f"ch{si}.case{ci}", "title": title, "degraded": notes})
    return drops


def print_fidelity(drops: list[dict], n_slides: int, deck_notes: list[str],
                   as_json=False, label="slides exported"):
    if as_json:
        print(json.dumps({"slides": n_slides, "degraded": drops,
                          "deck": deck_notes}, indent=2))
        return
    print(f"\n── Fidelity report ── {n_slides} {label}")
    if not drops:
        print("  nothing degraded per-slide — no advanced layouts/cyclers/iframes/videos in these chapters")
    for d in drops:
        print(f"  {d['slide']:>10}  {d['title'][:48]}")
        for n in d["degraded"]:
            print(f"            · {n}")
    for n in deck_notes:
        print(f"  ({n})")


STRUCTURAL_DECK_NOTES = [
    "structural mode is an editable two-column approximation — run with --visual for pixel fidelity",
    "cover/bonus/map/close slides not exported (SECTIONS chapters+cases only)",
    "custom theme reduced to 4 fixed accent RGBs on a dark background",
]


# ── Visual mode (pixel-faithful screenshots via capture_slides.py) ──────────

_TAG_RE = re.compile(r"<[^>]+>")


def _norm_title(s: str, drop_space=False) -> str:
    s = _TAG_RE.sub("", s or "").replace("\n", " ")
    s = re.sub(r"\s+", "" if drop_space else " ", s).strip().lower()
    return s


def _titles_match(a: str, b: str) -> bool:
    """Tolerant compare: whitespace-collapsed, tag-stripped, prefix-friendly
    (the manifest truncates titles at 120 chars; <br> joins can eat spaces)."""
    na, nb = _norm_title(a), _norm_title(b)
    if na == nb or na.startswith(nb) or nb.startswith(na):
        return True
    xa, xb = _norm_title(a, drop_space=True), _norm_title(b, drop_space=True)
    return xa == xb or xa.startswith(xb) or xb.startswith(xa)


# Mirrors getSlideNotes() defaults in index.html for slides outside SECTIONS.
DEFAULT_NOTES = {
    "cover": "Welcome the audience. Wait for stragglers. Introduce yourself and the talk theme.",
    "bonus": "Bonus lesson: rally the audience. This is the call to action. Build energy and end strong.",
    "map": "The constellation map is a visual summary of the journey. Let it breathe for a moment before moving on.",
    "close": "Wrap up. Point to the QR codes. Thank the audience. Open for questions.",
}


def build_notes_map(sections: list[dict], manifest: dict) -> dict[int, str]:
    """Map absolute slide index -> speaker notes, by aligning the manifest's
    type sequence with the deck's build order (settings, cover, then
    lesson+cases per chapter, then bonus/map/close). Lesson/case titles are
    cross-checked against SECTIONS; on any mismatch the slide gets NO notes —
    never wrong notes."""
    expected: list[tuple[str, tuple[int, int] | None]] = [("settings", None), ("cover", None)]
    for si, ch in enumerate(sections):
        expected.append(("lesson", (si, -1)))
        for ci in range(len(ch.get("cases") or [])):
            expected.append(("case", (si, ci)))
    expected += [("bonus", None), ("map", None), ("close", None)]

    slides = manifest.get("slides") or []
    aligned = len(slides) == len(expected) and all(
        (s.get("type") or "") == e[0] for s, e in zip(slides, expected))

    notes: dict[int, str] = {}
    for pos, s in enumerate(slides):
        ty = s.get("type") or ""
        if ty in DEFAULT_NOTES:
            notes[s["i"]] = DEFAULT_NOTES[ty]
            continue
        if not aligned or ty not in ("lesson", "case"):
            continue
        si, ci = expected[pos][1]  # type: ignore[misc]
        if ty == "lesson":
            src = sections[si].get("lesson") or {}
        else:
            src = (sections[si].get("cases") or [])[ci]
        title = src.get("title") or ""
        if title and (s.get("title") or "") and not _titles_match(title, s["title"]):
            continue  # ambiguous alignment — attach nothing rather than wrong notes
        txt = (src.get("notes") or "").strip()
        if txt:
            notes[s["i"]] = txt
    return notes


def visual_fidelity_report(manifest: dict, exported_indices: list[int], as_json=False):
    """Manifest-based report (mirrors export_pdf.py): what the deck's own
    export mode froze or swapped while rendering the screenshots."""
    drops = []
    for s in manifest["slides"]:
        if s["i"] not in exported_indices:
            continue
        notes = []
        if s.get("cycler", 0) > 1:
            notes.append(f"media cycler: first of {s['cycler']} items shown")
        if s.get("iframes"):
            notes.append(f"{s['iframes']} iframe(s) -> placeholder panel with URL")
        if s.get("videos"):
            notes.append(f"{s['videos']} video(s) frozen on an early frame")
        if s["type"] == "map":
            notes.append("constellation map: 3D/animated content captured as-rendered")
        if notes:
            drops.append({"slide": s["i"], "title": s["title"], "degraded": notes})
    if as_json:
        print(json.dumps({"slides": len(exported_indices), "degraded": drops,
                          "deck": ["visual mode: slides are screenshots — pixel-faithful but not editable"]},
                         indent=2))
        return
    print(f"\n── Fidelity report ── {len(exported_indices)} slides exported (visual)")
    if not drops:
        print("  nothing degraded — interactive-only features were not used on these slides")
    for d in drops:
        print(f"  slide {d['slide']:>3}  {d['title'][:48]}")
        for n in d["degraded"]:
            print(f"            · {n}")
    hidden = [s for s in manifest["slides"] if s["hidden"] and s["type"] != "settings"]
    if hidden:
        print(f"  ({len(hidden)} hidden slide(s) excluded, as in the live deck)")
    print("  (visual mode: slides are screenshots — pixel-faithful but not editable)")


def build_visual(args) -> int:
    from capture_slides import find_chrome, read_manifest, shoot_all, parse_slide_spec
    chrome = find_chrome()  # sys.exits with an actionable message if absent
    print("reading deck manifest…", file=sys.stderr)
    man = read_manifest(str(args.html), chrome)
    indices = sorted(parse_slide_spec(args.slides, man["visible"]))

    shot_dir = args.keep_shots or tempfile.mkdtemp(prefix="sd-pptx-shots-")
    print(f"capturing {len(indices)} slides at {args.width}x{args.height}…", file=sys.stderr)
    shots = shoot_all(str(args.html), indices, shot_dir, args.width, args.height,
                      chrome, args.parallel)

    sections = extract_sections(args.html)
    notes_map = build_notes_map(sections, man)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for i in indices:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _fill_bg(slide)
        slide.shapes.add_picture(shots[i], 0, 0, width=SLIDE_W, height=SLIDE_H)
        txt = notes_map.get(i)
        if txt:
            slide.notes_slide.notes_text_frame.text = txt
    prs.save(str(args.out))
    if not args.keep_shots:
        shutil.rmtree(shot_dir, ignore_errors=True)

    size_mb = Path(args.out).stat().st_size / 1e6
    print(f"[done] Wrote {args.out} — {len(indices)} slides (visual, pixel-faithful), "
          f"{size_mb:.1f} MB", file=sys.stderr)
    if not args.no_report:
        visual_fidelity_report(man, indices, args.json)
    return 0


def main() -> int:
    ap = argparse.ArgumentParser(description="Export the deck as .pptx — "
                                 "structural (editable, default) or --visual (pixel-faithful)")
    ap.add_argument("--html", type=Path, default=DEFAULT_HTML)
    ap.add_argument("--out", type=Path, required=True)
    ap.add_argument("--chapter", type=int, default=None,
                    help="structural mode: export a single chapter")
    ap.add_argument("--visual", action="store_true",
                    help="pixel-faithful screenshot slides via headless Chrome (not editable)")
    ap.add_argument("--slides", help="--visual only: absolute indices, e.g. '1,3-5' (default: all visible)")
    ap.add_argument("--width", type=int, default=1920, help="--visual capture width")
    ap.add_argument("--height", type=int, default=1080, help="--visual capture height")
    ap.add_argument("--parallel", type=int, default=4, help="--visual concurrent captures")
    ap.add_argument("--keep-shots", metavar="DIR", help="--visual: keep per-slide PNGs here")
    ap.add_argument("--json", action="store_true", help="fidelity report as JSON")
    ap.add_argument("--no-report", action="store_true", help="skip the fidelity report")
    args = ap.parse_args()

    if args.visual:
        if args.chapter is not None:
            print("ERROR: --chapter is structural-only; use --slides with --visual", file=sys.stderr)
            return 2
        return build_visual(args)

    sections = extract_sections(args.html)
    if args.chapter is not None:
        if not (0 <= args.chapter < len(sections)):
            print(f"ERROR: chapter index out of range (0..{len(sections)-1})", file=sys.stderr)
            return 2
        sections = [sections[args.chapter]]

    n_ch, n_cases = build(sections, args.out)
    print(f"[done] Wrote {args.out} ({n_ch} chapters, {n_cases} cases, {n_ch + n_cases} slides)", file=sys.stderr)
    if not args.no_report:
        print_fidelity(structural_fidelity(sections), n_ch + n_cases,
                       STRUCTURAL_DECK_NOTES, args.json,
                       label="slides exported (structural)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
